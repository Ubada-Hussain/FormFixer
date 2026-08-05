from fastapi import FastAPI, File, UploadFile, HTTPException, Header
from fastapi.responses import Response
import tempfile
import os
import uuid
import shutil
import logging

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = FastAPI(title="pdf2docx-service", version="1.0.0")


@app.get("/health")
def health():
    return {"status": "up"}


MAX_FILE_SIZE = 10 * 1024 * 1024  # 10 MB
MAX_PAGES = 200


@app.post("/convert")
async def convert_pdf_to_docx(
    file: UploadFile = File(...),
    x_internal_key: str = Header(None)
):
    expected_key = os.environ.get("INTERNAL_API_KEY")
    if expected_key and x_internal_key != expected_key:
        raise HTTPException(status_code=401, detail="Unauthorized: Invalid internal API key.")

    if not file.filename:
        raise HTTPException(status_code=400, detail="No filename provided.")

    if not file.filename.lower().endswith(".pdf"):
        raise HTTPException(
            status_code=422,
            detail=f"Unsupported file type: '{file.filename}'. Only PDF files are accepted.",
        )

    pdf_bytes = await file.read()
    if len(pdf_bytes) == 0:
        raise HTTPException(status_code=400, detail="Uploaded file is empty.")

    if len(pdf_bytes) > MAX_FILE_SIZE:
        raise HTTPException(
            status_code=413,
            detail=f"File is too large ({len(pdf_bytes) / 1024 / 1024:.1f} MB). Maximum allowed size is 10 MB.",
        )

    logger.info(f"Converting '{file.filename}' ({len(pdf_bytes):,} bytes) → .docx")

    tmp_dir = tempfile.mkdtemp()
    try:
        pdf_path = os.path.join(tmp_dir, f"{uuid.uuid4()}.pdf")
        docx_path = pdf_path.replace(".pdf", ".docx")

        with open(pdf_path, "wb") as f:
            f.write(pdf_bytes)

        # 1. RTL / Urdu / Arabic Check + page count guard
        import fitz
        doc = fitz.open(pdf_path)

        if len(doc) > MAX_PAGES:
            doc.close()
            raise HTTPException(
                status_code=413,
                detail=f"PDF has {len(doc)} pages. Maximum allowed is {MAX_PAGES} pages.",
            )

        sample_text = ""
        for i in range(min(2, len(doc))):
            sample_text += doc[i].get_text()
        doc.close()
        
        rtl_count = sum(1 for c in sample_text if '\u0600' <= c <= '\u06FF')
        if len(sample_text) > 50 and (rtl_count / len(sample_text)) > 0.03:
            return Response(
                status_code=400,
                content='{"error": "Urdu/Arabic PDF-to-Word conversion isn\'t supported yet — we don\'t want to hand you garbled output. This is coming soon. English documents convert normally.", "isRtlUnsupported": true}',
                media_type="application/json"
            )

        # 2. Proceed with conversion if no RTL
        # Imported here to prevent heavy load at startup
        from pdf2docx import Converter
        cv = Converter(pdf_path)
        cv.convert(docx_path, start=0, end=None)
        cv.close()

        if not os.path.exists(docx_path):
            raise HTTPException(
                status_code=500,
                detail="Conversion produced no output file — the PDF may be image-only or encrypted.",
            )

        with open(docx_path, "rb") as f:
            docx_bytes = f.read()

        output_name = os.path.splitext(file.filename)[0] + "-converted.docx"
        logger.info(f"Conversion OK → '{output_name}' ({len(docx_bytes):,} bytes)")

        return Response(
            content=docx_bytes,
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={
                "Content-Disposition": f'attachment; filename="{output_name}"',
                "Content-Length": str(len(docx_bytes)),
            },
        )

    except HTTPException:
        raise
    except Exception as e:
        # Log the full traceback server-side, but never expose it to the caller
        logger.error(f"Conversion error: {e}", exc_info=True)
        raise HTTPException(
            status_code=500,
            detail="Conversion failed due to an internal error. Please try a different file.",
        )
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)

@app.post("/pdf-to-pptx")
async def convert_pdf_to_pptx(
    file: UploadFile = File(...),
    x_internal_key: str = Header(None)
):
    expected_key = os.environ.get("INTERNAL_API_KEY")
    if expected_key and x_internal_key != expected_key:
        raise HTTPException(status_code=401, detail="Unauthorized: Invalid internal API key.")

    if not file.filename:
        raise HTTPException(status_code=400, detail="No filename provided.")

    if not file.filename.lower().endswith(".pdf"):
        raise HTTPException(
            status_code=422,
            detail=f"Unsupported file type: '{file.filename}'. Only PDF files are accepted.",
        )

    pdf_bytes = await file.read()
    if len(pdf_bytes) == 0:
        raise HTTPException(status_code=400, detail="Uploaded file is empty.")

    if len(pdf_bytes) > MAX_FILE_SIZE:
        raise HTTPException(
            status_code=413,
            detail=f"File is too large ({len(pdf_bytes) / 1024 / 1024:.1f} MB). Maximum allowed size is {MAX_FILE_SIZE / 1024 / 1024:.0f} MB.",
        )

    logger.info(f"Converting '{file.filename}' ({len(pdf_bytes):,} bytes) → .pptx")

    tmp_dir = tempfile.mkdtemp()
    try:
        pdf_path = os.path.join(tmp_dir, f"{uuid.uuid4()}.pdf")
        pptx_path = pdf_path.replace(".pdf", ".pptx")

        with open(pdf_path, "wb") as f:
            f.write(pdf_bytes)

        # Page count guard
        import fitz
        doc = fitz.open(pdf_path)
        page_count = len(doc)
        doc.close()
        
        if page_count > MAX_PAGES:
            raise HTTPException(
                status_code=413,
                detail=f"PDF has {page_count} pages. Maximum allowed is {MAX_PAGES} pages.",
            )

        from pdf2image import convert_from_path
        from pptx import Presentation
        from pptx.util import Inches

        # Render PDF pages to images
        images = convert_from_path(pdf_path, dpi=300)
        
        if not images:
            raise HTTPException(status_code=500, detail="Failed to render any pages from the PDF.")

        prs = Presentation()
        # default pptx slide width is 10 inches, height is 7.5 inches (4:3) or 13.33 x 7.5 (16:9).
        # We will match the aspect ratio of the first image.
        first_img = images[0]
        img_width, img_height = first_img.size
        
        # Keep width at 10 inches and adjust height based on aspect ratio
        base_width = Inches(10)
        aspect_ratio = img_height / img_width
        base_height = int(base_width * aspect_ratio)

        prs.slide_width = base_width
        prs.slide_height = base_height

        blank_slide_layout = prs.slide_layouts[6]

        for i, image in enumerate(images):
            img_path = os.path.join(tmp_dir, f"page_{i}.png")
            image.save(img_path, "PNG")
            
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(img_path, 0, 0, prs.slide_width, prs.slide_height)

        prs.save(pptx_path)

        if not os.path.exists(pptx_path):
            raise HTTPException(
                status_code=500,
                detail="Conversion produced no output file.",
            )

        with open(pptx_path, "rb") as f:
            pptx_bytes = f.read()

        output_name = os.path.splitext(file.filename)[0] + "-converted.pptx"
        logger.info(f"Conversion OK → '{output_name}' ({len(pptx_bytes):,} bytes)")

        return Response(
            content=pptx_bytes,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={
                "Content-Disposition": f'attachment; filename="{output_name}"',
                "Content-Length": str(len(pptx_bytes)),
            },
        )

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Conversion error: {e}", exc_info=True)
        raise HTTPException(
            status_code=500,
            detail="Conversion failed due to an internal error. Please try a different file.",
        )
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)

if __name__ == "__main__":
    import uvicorn
    # Railway passes the PORT environment variable.
    # Programmatic parsing guarantees it binds correctly regardless of Docker CMD shell expansion bugs.
    port = int(os.environ.get("PORT", 8000))
    logger.info(f"Starting uvicorn on port {port}")
    uvicorn.run(app, host="0.0.0.0", port=port)
